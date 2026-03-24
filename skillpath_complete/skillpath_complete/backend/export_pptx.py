from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io

BG    = RGBColor(0x12, 0x12, 0x12)
CARD  = RGBColor(0x1e, 0x1e, 0x24)
WHITE = RGBColor(0xf3, 0xf4, 0xf6)
MUTED = RGBColor(0x9c, 0xa3, 0xaf)
BLUE  = RGBColor(0x63, 0x66, 0xf1)
CYAN  = RGBColor(0x22, 0xd3, 0xee)
GREEN = RGBColor(0x10, 0xb9, 0x81)
AMBER = RGBColor(0xf5, 0x9e, 0x0b)
RED   = RGBColor(0xef, 0x44, 0x44)

def hex_to_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def set_bg(slide, prs, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height,
                font_size=14, bold=False, color=None,
                align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or WHITE
    return txBox

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def build_pptx(data: dict) -> bytes:
    resume  = data.get("resume", {})
    job     = data.get("job", {})
    pathway = data.get("pathway", {})
    diag    = data.get("diagnostic")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # SLIDE 1: TITLE
    s1 = prs.slides.add_slide(blank)
    set_bg(s1, prs, BG)
    add_rect(s1, 0, 0, 0.08, 7.5, BLUE)
    add_rect(s1, 0, 0, 13.33, 0.06, BLUE)
    add_textbox(s1, "SKILLPATH", 0.3, 0.3, 8, 0.5, font_size=11, bold=True, color=BLUE)
    add_textbox(s1, "AI-Adaptive Onboarding Engine", 0.3, 0.7, 12, 0.5, font_size=11, color=MUTED)
    add_textbox(s1, "Personalised Learning Pathway", 0.3, 1.4, 12, 0.9, font_size=36, bold=True, color=WHITE)
    name = resume.get("name", "Candidate")
    role = job.get("role_title", "Target Role")
    seniority = job.get("seniority", "")
    add_textbox(s1, f"{name}  →  {role} ({seniority})", 0.3, 2.5, 12, 0.5, font_size=14, color=CYAN)
    summary = pathway.get("skill_gap_summary", "")
    if summary:
        add_textbox(s1, summary, 0.3, 3.2, 11, 1.2, font_size=11, color=MUTED, italic=True)
    score = pathway.get("readiness_score", 0)
    days  = pathway.get("estimated_total_days", 0)
    saved = pathway.get("time_saved_days", 0)
    stats = [
        (str(score)+"%", "Readiness Score",   BLUE),
        (str(days)+"d",  "Training Duration", CYAN),
        (str(saved)+"d", "Days Saved",        GREEN),
    ]
    for i, (val, lbl, col) in enumerate(stats):
        x = 0.3 + i * 2.6
        add_rect(s1, x, 4.7, 2.3, 1.2, CARD, col)
        add_textbox(s1, val, x+0.1, 4.85, 2.1, 0.55, font_size=22, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(s1, lbl, x+0.1, 5.4,  2.1, 0.35, font_size=9,  color=MUTED, align=PP_ALIGN.CENTER)
    mode = "Diagnostic Mode" if diag else "Resume Mode"
    add_textbox(s1, mode, 10, 4.7, 3, 0.4, font_size=10, color=BLUE, bold=True, align=PP_ALIGN.RIGHT)
    if diag:
        add_textbox(s1,
            f"Quiz Score: {diag.get('score_pct',0)}%  |  {diag.get('total_questions',0)} questions",
            10, 5.15, 3, 0.4, font_size=9, color=MUTED, align=PP_ALIGN.RIGHT)

    # SLIDE 2: SKILL ANALYSIS
    s2 = prs.slides.add_slide(blank)
    set_bg(s2, prs, BG)
    add_rect(s2, 0, 0, 0.08, 7.5, CYAN)
    add_textbox(s2, "SKILL ANALYSIS", 0.3, 0.2, 12, 0.5, font_size=11, bold=True, color=MUTED)
    add_textbox(s2, "Where You Stand vs Role Requirements", 0.3, 0.6, 12, 0.6, font_size=24, bold=True, color=WHITE)
    present = pathway.get("skills_present", [])
    partial = pathway.get("skills_partial", [])
    missing = pathway.get("skills_missing", [])
    cols = [
        (present, "SKILLS VERIFIED", GREEN, 0.3),
        (partial, "NEEDS DEEPENING", AMBER, 4.5),
        (missing, "GAPS TO CLOSE",   RED,   8.7),
    ]
    for skills, title, col, x in cols:
        add_rect(s2, x, 1.5, 3.8, 0.45, CARD, col)
        add_textbox(s2, f"{title}  ({len(skills)})", x+0.12, 1.55, 3.6, 0.35, font_size=10, bold=True, color=col)
        y = 2.1
        for sk in skills[:8]:
            add_rect(s2, x, y, 3.8, 0.3, CARD)
            add_textbox(s2, "• " + sk, x+0.12, y+0.04, 3.5, 0.25, font_size=9, color=WHITE)
            y += 0.33
        if len(skills) > 8:
            add_textbox(s2, f"+ {len(skills)-8} more", x+0.12, y+0.04, 3.5, 0.25, font_size=8, color=MUTED)

    # SLIDES 3+: LEARNING PATHWAY CARDS
    steps = pathway.get("pathway", [])
    slide_groups = [steps[i:i+3] for i in range(0, len(steps), 3)]
    for gi, group in enumerate(slide_groups):
        s = prs.slides.add_slide(blank)
        set_bg(s, prs, BG)
        add_rect(s, 0, 0, 0.08, 7.5, GREEN)
        add_textbox(s, "LEARNING PATHWAY", 0.3, 0.2, 12, 0.4, font_size=11, bold=True, color=MUTED)
        add_textbox(s,
            f"Personalised Courses  ·  {pathway.get('estimated_total_days',0)} Days Total  ·  {len(steps)} Modules",
            0.3, 0.55, 12, 0.4, font_size=20, bold=True, color=WHITE)
        level_colors = {"beginner": GREEN, "intermediate": AMBER, "advanced": RED}
        card_w = 4.1
        for ci, step in enumerate(group):
            x = 0.25 + ci * (card_w + 0.15)
            col = level_colors.get(step.get("level", "beginner"), BLUE)
            add_rect(s, x, 1.2, card_w, 5.8, CARD, col)
            add_rect(s, x, 1.2, card_w, 0.45, col)
            add_textbox(s,
                f"#{step.get('order','')}  {step.get('level','').upper()}  ·  {step.get('duration','')}",
                x+0.12, 1.22, card_w-0.2, 0.38, font_size=9, bold=True, color=BG)
            add_textbox(s, step.get("course_title",""), x+0.12, 1.75, card_w-0.2, 0.7, font_size=12, bold=True, color=WHITE)
            desc = step.get("description","")
            if desc:
                add_textbox(s, desc, x+0.12, 2.55, card_w-0.2, 0.7, font_size=8, color=MUTED, italic=True)
            add_rect(s, x+0.08, 3.35, card_w-0.16, 0.28, hex_to_rgb("#0f0f1a"))
            add_textbox(s, "AI REASONING", x+0.15, 3.37, card_w-0.3, 0.22, font_size=7, bold=True, color=col)
            add_textbox(s, step.get("reasoning",""), x+0.12, 3.68, card_w-0.2, 1.4, font_size=8, color=WHITE)
            outcome = step.get("expected_outcome","")
            if outcome:
                add_rect(s, x+0.08, 5.18, card_w-0.16, 0.22, CARD)
                add_textbox(s, "OUTCOME: " + outcome, x+0.12, 5.2, card_w-0.2, 0.55, font_size=7.5, color=CYAN)
            gaps = step.get("addresses_skill_gap",[])
            if gaps:
                add_textbox(s, "Closes: " + ", ".join(gaps[:3]), x+0.12, 5.85, card_w-0.2, 0.35, font_size=7, color=col)

    # FINAL SLIDE: REASONING TRACE
    sf = prs.slides.add_slide(blank)
    set_bg(sf, prs, BG)
    add_rect(sf, 0, 0, 0.08, 7.5, AMBER)
    add_textbox(sf, "AI REASONING TRACE", 0.3, 0.2, 12, 0.4, font_size=11, bold=True, color=MUTED)
    add_textbox(sf, "How and Why This Pathway Was Built", 0.3, 0.55, 12, 0.5, font_size=22, bold=True, color=WHITE)
    trace = pathway.get("reasoning_trace", {})
    if trace:
        add_rect(sf, 0.3, 1.25, 12.73, 0.9, CARD, AMBER)
        add_textbox(sf, "SEQUENCING APPROACH", 0.45, 1.3, 4, 0.3, font_size=8, bold=True, color=AMBER)
        add_textbox(sf, trace.get("approach","—"), 0.45, 1.58, 12.4, 0.5, font_size=9, color=WHITE)
        conf = trace.get("confidence","").upper()
        conf_col = GREEN if conf=="HIGH" else AMBER if conf=="MEDIUM" else RED
        add_rect(sf, 0.3, 2.35, 6.2, 0.35, CARD, conf_col)
        add_textbox(sf, f"CONFIDENCE: {conf}  ·  {trace.get('confidence_reason','')}", 0.45, 2.38, 6, 0.28, font_size=9, color=conf_col)
        prereqs = trace.get("prerequisites_applied", [])
        add_rect(sf, 0.3, 2.95, 6.0, 0.32, CARD)
        add_textbox(sf, "PREREQUISITE CHAINS APPLIED", 0.45, 2.97, 5.7, 0.26, font_size=8, bold=True, color=BLUE)
        for i, p in enumerate(prereqs[:5]):
            add_textbox(sf, "→ " + p, 0.45, 3.3 + i*0.3, 5.7, 0.28, font_size=8.5, color=WHITE)
        skipped = trace.get("skipped_courses", [])
        add_rect(sf, 6.8, 2.95, 6.23, 0.32, CARD)
        add_textbox(sf, f"COURSES SKIPPED ({len(skipped)}) — ADAPTIVE LOGIC", 6.95, 2.97, 6, 0.26, font_size=8, bold=True, color=RED)
        for i, sk in enumerate(skipped[:5]):
            add_textbox(sf, "✕ " + sk, 6.95, 3.3 + i*0.3, 6, 0.28, font_size=8.5, color=MUTED)
    add_textbox(sf, "Generated by SkillPath AI · Powered by Claude",
        0.3, 6.9, 12.73, 0.35, font_size=9, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

